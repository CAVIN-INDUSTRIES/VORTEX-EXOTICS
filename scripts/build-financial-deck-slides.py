#!/usr/bin/env python3
"""Append financial visualization slides to VEX-Pitch-Deck.pptx.

Real figures: tier list prices from apps/api/src/routes/pricing.ts (PLANS).
Illustrative charts: clearly labeled; replace with finance-approved data before
external distribution.
"""

from __future__ import annotations

import shutil
import sys
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches, Pt

REPO = Path(__file__).resolve().parents[1]
PACKAGE = REPO / "VEX-Business-Package"
PPTX = PACKAGE / "VEX-Pitch-Deck.pptx"

# Source: apps/api/src/routes/pricing.ts PLANS
TIERS = ("Starter", "Pro", "Enterprise")
MONTHLY = (49, 149, 299)
YEARLY = (470, 1430, 2870)
MONTHLY_X12 = tuple(m * 12 for m in MONTHLY)

DISCLAIMER = (
    "Illustrative projection — replace with finance / Stripe-backed figures before external use."
)
SOURCE_NOTE = "List prices: product API (apps/api/src/routes/pricing.ts)."


def _add_footer(slide, text: str, top: float = 6.85) -> None:
    box = slide.shapes.add_textbox(Inches(0.45), Inches(top), Inches(12.4), Inches(0.55))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.italic = True
    p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)


def _chart_title(chart, title: str) -> None:
    chart.has_title = True
    chart.chart_title.text_frame.text = title
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)


def main() -> int:
    if not PPTX.is_file():
        print("Missing", PPTX, file=sys.stderr)
        return 1

    backup = Path(tempfile.gettempdir()) / "VEX-Pitch-Deck-backup.pptx"
    shutil.copy2(PPTX, backup)

    prs = Presentation(PPTX)
    # 16:9 deck — use Title Only (5) for chart slides, Section Header (2) for breaks

    # --- Section divider
    s = prs.slides.add_slide(prs.slide_layouts[2])
    s.shapes.title.text = "Financial model & projections"
    st = s.placeholders[1]
    st.text = "Source-backed pricing · Illustrative scenarios (finance to validate)\n" + SOURCE_NOTE

    # --- 1) Monthly list price (REAL)
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "Published SaaS list prices (monthly)"
    data = CategoryChartData()
    data.categories = TIERS
    data.add_series("USD / mo", MONTHLY)
    x, y, cx, cy = Inches(0.55), Inches(1.35), Inches(12.3), Inches(5.2)
    chart = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, data).chart
    _chart_title(chart, "")
    chart.value_axis.has_major_gridlines = True
    _add_footer(s, SOURCE_NOTE, 6.78)

    # --- 2) Yearly list price (REAL)
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "Published SaaS list prices (yearly contract)"
    data = CategoryChartData()
    data.categories = TIERS
    data.add_series("USD / yr", YEARLY)
    chart = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.55), Inches(1.35), Inches(12.3), Inches(5.2), data).chart
    chart.value_axis.has_major_gridlines = True
    _add_footer(s, SOURCE_NOTE, 6.78)

    # --- 3) Monthly×12 vs yearly (REAL savings signal)
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "Yearly list vs paying monthly × 12"
    data = CategoryChartData()
    data.categories = TIERS
    data.add_series("12 × monthly list", MONTHLY_X12)
    data.add_series("Yearly list", YEARLY)
    chart = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.55), Inches(1.35), Inches(12.3), Inches(5.2), data).chart
    if chart.legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.value_axis.has_major_gridlines = True
    _add_footer(s, "Shows list-price delta; actual invoice depends on Stripe Price IDs.", 6.78)

    # --- 4) Illustrative ARR paths (3 scenarios)
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "Illustrative total scale ($M) — bear / base / bull"
    years = ("Y1", "Y2", "Y3", "Y4", "Y5")
    data = CategoryChartData()
    data.categories = years
    data.add_series("Bear", (0.42, 1.35, 3.12, 5.72, 9.45))
    data.add_series("Base", (0.50, 1.67, 3.57, 6.12, 12.50))
    data.add_series("Bull", (0.55, 1.84, 3.94, 6.74, 13.77))
    chart = s.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, Inches(0.55), Inches(1.35), Inches(12.3), Inches(5.2), data).chart
    if chart.legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    _chart_title(chart, "")
    chart.value_axis.has_major_gridlines = True
    _add_footer(s, DISCLAIMER, 6.72)

    # --- 5) Y5 scenario column
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "Illustrative year-5 ARR by scenario ($M)"
    data = CategoryChartData()
    data.categories = ("Bear", "Base", "Bull")
    data.add_series("ARR $M", (9.45, 12.50, 13.77))
    chart = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.55), Inches(1.35), Inches(12.3), Inches(5.2), data).chart
    chart.value_axis.has_major_gridlines = True
    _add_footer(s, DISCLAIMER, 6.72)

    # --- 6) Pie — illustrative revenue mix Y5 base
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "Illustrative revenue mix (year 5, base scenario)"
    data = CategoryChartData()
    data.categories = (
        "Subscription\n(Starter/Pro)",
        "Enterprise tier",
        "Success / transaction",
        "Add-ons / services",
    )
    data.add_series("Mix %", (62, 18, 12, 8))
    chart = s.shapes.add_chart(XL_CHART_TYPE.PIE, Inches(0.55), Inches(1.35), Inches(12.3), Inches(5.2), data).chart
    chart.plots[0].has_data_labels = True
    _add_footer(s, DISCLAIMER + " Percentages are planning placeholders.", 6.72)

    # --- 7) Stacked column — subscription vs usage vs success (illustrative $M)
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "Illustrative revenue build ($M ARR eq.) — stacked components"
    data = CategoryChartData()
    data.categories = years
    data.add_series("Subscription + Ultra", (0.32, 0.99, 1.98, 3.15, 4.37))
    data.add_series("Success fees (ann.)", (0.09, 0.38, 0.95, 1.85, 2.10))
    chart = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED, Inches(0.55), Inches(1.35), Inches(12.3), Inches(5.2), data).chart
    if chart.legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.value_axis.has_major_gridlines = True
    _add_footer(s, DISCLAIMER + " Pre-uplift components; not GAAP.", 6.68)

    # --- 8) Illustrative EBITDA ($M)
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "Illustrative EBITDA path ($M) — base scenario"
    data = CategoryChartData()
    data.categories = years
    data.add_series("EBITDA", (-1.8, -0.9, 0.2, 0.85, 1.35))
    chart = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.55), Inches(1.35), Inches(12.3), Inches(5.2), data).chart
    chart.value_axis.has_major_gridlines = True
    _add_footer(s, DISCLAIMER, 6.72)

    # --- 9) Opex allocation pie (planning illustration)
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "Illustrative mature opex allocation (% of revenue)"
    data = CategoryChartData()
    data.categories = ("R&D / product", "Sales & marketing", "G&A", "Infra", "Support")
    data.add_series("%", (38, 28, 14, 12, 8))
    chart = s.shapes.add_chart(XL_CHART_TYPE.PIE, Inches(0.55), Inches(1.35), Inches(12.3), Inches(5.2), data).chart
    chart.plots[0].has_data_labels = True
    _add_footer(s, "Planning illustration only — not actual cost accounting.", 6.72)

    # --- 10) Sensitivity heatmap substitute: clustered bar for ARPU shock
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "Illustrative Y5 ARR sensitivity — blended ARPU shock"
    data = CategoryChartData()
    data.categories = ("ARPU −15%", "Base ARPU", "ARPU +15%")
    data.add_series("$M", (10.6, 12.5, 14.4))
    chart = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.55), Inches(1.35), Inches(12.3), Inches(5.2), data).chart
    chart.value_axis.has_major_gridlines = True
    _add_footer(s, DISCLAIMER, 6.72)

    # --- 11) Waterfall substitute: bar chart steps (illustrative MRR bridge)
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "Illustrative MRR bridge ($k / mo) — template"
    data = CategoryChartData()
    data.categories = ("Start", "+ New logos", "+ Expansion", "− Churn", "= End")
    data.add_series("$k/mo", (45, 28, 15, -8, 80))
    chart = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.55), Inches(1.35), Inches(12.3), Inches(5.2), data).chart
    chart.value_axis.has_major_gridlines = True
    _add_footer(s, DISCLAIMER + " Numeric example for chart style only.", 6.68)

    # --- 12) Area — cumulative cash feel (illustrative)
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "Illustrative cumulative revenue index (base = 100 at Y1)"
    data = CategoryChartData()
    data.categories = years
    data.add_series("Index", (100, 334, 714, 1224, 2500))
    chart = s.shapes.add_chart(XL_CHART_TYPE.AREA, Inches(0.55), Inches(1.35), Inches(12.3), Inches(5.2), data).chart
    chart.value_axis.has_major_gridlines = True
    _add_footer(s, DISCLAIMER + " Indexed shape for narrative only.", 6.72)

    # --- 13) Text: data sources
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "Financial charts — data sources & next steps"
    lines = (
        "REAL (use as-is): Monthly $49 / $149 / $299 and yearly $470 / $1,430 / $2,870 from GET /pricing/plans (pricing.ts).",
        "ILLUSTRATIVE: All scenario curves, EBITDA, mix %, sensitivity, and MRR bridge slides — replace with FP&A model tied to Stripe.",
        "Live metrics: RaisePackage + pilot aggregates (capital routes, pilotMetrics.ts).",
        "Before investors: delete or refresh illustrative slides; keep source-backed pricing slides.",
    )
    tf = s.placeholders[1].text_frame
    tf.text = lines[0]
    for line in lines[1:]:
        p = tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(14)
    tf.paragraphs[0].font.size = Pt(14)

    prs.save(PPTX)
    print("Updated", PPTX, "slides:", len(prs.slides), "backup:", backup)
    print("Next: python scripts/apply-vex-pitch-deck-theme.py  (obsidian/violet theme + public-friendly titles)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
