#!/usr/bin/env python3
"""
Apply VEX obsidian / violet–gold visual theme + accessible financial verbiage
to VEX-Business-Package/VEX-Pitch-Deck.pptx.

Run after build-financial-deck-slides.py (or any time the deck drifts).
Uses /tmp backup before overwrite.
"""

from __future__ import annotations

import shutil
import sys
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.util import Pt

REPO = Path(__file__).resolve().parents[1]
PPTX = REPO / "VEX-Business-Package" / "VEX-Pitch-Deck.pptx"

# VEX luxury / futuristic palette (obsidian vault + violet + gold + cool neutrals)
OBSIDIAN = RGBColor(0x0A, 0x0A, 0x12)
OBSIDIAN_ELEV = RGBColor(0x12, 0x12, 0x1E)
PLATINUM = RGBColor(0xF0, 0xEE, 0xF8)
MUTED = RGBColor(0xA8, 0xA4, 0xBC)
VIOLET = RGBColor(0xA7, 0x8B, 0xFA)  # soft violet
VIOLET_DEEP = RGBColor(0x7C, 0x3A, 0xED)
GOLD = RGBColor(0xE8, 0xC5, 0x47)
GOLD_SOFT = RGBColor(0xD4, 0xAF, 0x37)
ACCENT_CYAN = RGBColor(0x5E, 0xD4, 0xE0)
GRID = RGBColor(0x38, 0x38, 0x50)
CARD_FILL = RGBColor(0x16, 0x16, 0x24)
CARD_LINE = RGBColor(0x4C, 0x3D, 0x6B)

FONT_UI = "Segoe UI"
FONT_FALLBACK = "Calibri"

# Series colors for charts (rotate)
CHART_COLORS = (
    VIOLET,
    GOLD,
    ACCENT_CYAN,
    RGBColor(0xF4, 0x7B, 0xA3),
    PLATINUM,
    MUTED,
)


def delete_slide(prs: Presentation, index: int) -> None:
    r_id = prs.slides._sldIdLst[index].rId
    prs.part.drop_rel(r_id)
    del prs.slides._sldIdLst[index]


def dedupe_financial_section(prs: Presentation) -> int:
    """Remove duplicate appended financial block (keeps first 32 slides if 46)."""
    removed = 0
    while len(prs.slides) > 32:
        delete_slide(prs, len(prs.slides) - 1)
        removed += 1
    return removed


def apply_slide_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = OBSIDIAN


def style_textframe(tf, *, title: bool = False, footer: bool = False) -> None:
    if not tf:
        return
    for p in tf.paragraphs:
        if footer:
            base = MUTED
            sz = Pt(9)
        elif title:
            base = PLATINUM
            sz = Pt(28) if p.level == 0 else Pt(20)
        else:
            base = PLATINUM
            sz = None
        for run in p.runs:
            run.font.name = FONT_UI
            run.font.color.rgb = base
            if sz and run.font.size is None:
                run.font.size = sz
        if not p.runs:
            p.font.name = FONT_UI
            p.font.color.rgb = base
            if sz:
                p.font.size = sz


def style_chart(chart) -> None:
    try:
        chart.font.size = Pt(11)
        chart.font.color.rgb = MUTED
    except Exception:
        pass
    try:
        chart.chart_title.text_frame.paragraphs[0].font.color.rgb = PLATINUM
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)
        chart.chart_title.text_frame.paragraphs[0].font.name = FONT_UI
    except Exception:
        pass
    for _axis_name in ("category_axis", "value_axis"):
        try:
            axis = getattr(chart, _axis_name)
        except ValueError:
            continue
        if axis is None:
            continue
        try:
            axis.tick_labels.font.color.rgb = MUTED
            axis.tick_labels.font.size = Pt(10)
            axis.tick_labels.font.name = FONT_UI
            axis.format.line.color.rgb = GRID
        except Exception:
            pass
    try:
        va = chart.value_axis
        va.has_major_gridlines = True
        va.major_gridlines.format.line.color.rgb = GRID
        va.major_gridlines.format.line.width = Pt(0.5)
    except Exception:
        pass
    try:
        if chart.legend:
            chart.legend.font.color.rgb = PLATINUM
            chart.legend.font.size = Pt(10)
            chart.legend.font.name = FONT_UI
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    except Exception:
        pass
    try:
        chart.plot_area.format.fill.solid()
        chart.plot_area.format.fill.fore_color.rgb = OBSIDIAN_ELEV
    except Exception:
        pass
    for idx, plot in enumerate(chart.plots):
        for sidx, series in enumerate(plot.series):
            try:
                f = series.format.fill
                f.solid()
                f.fore_color.rgb = CHART_COLORS[sidx % len(CHART_COLORS)]
            except Exception:
                pass


def recolor_auto_shapes(slide) -> None:
    """Subtly align deck chrome: dark cards + violet border hints."""
    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
            continue
        try:
            if shape.fill.type == MSO_FILL_TYPE.SOLID:
                # Skip huge full-slide rectangles if any
                if shape.width > 10_000_000:  # EMUs rough
                    continue
                rgb = shape.fill.fore_color.rgb
                if rgb is None:
                    continue
                r, g, b = rgb
                # Old deck: bright blues / teals / navy cards
                if b > r + 30 and b > g:  # blue-ish accent bars
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = VIOLET_DEEP
                elif r < 0x30 and g < 0x30 and b < 0x40:  # dark navy card
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = CARD_FILL
            if shape.line.fill.type == MSO_FILL_TYPE.SOLID:
                shape.line.color.rgb = CARD_LINE
        except Exception:
            continue


def replace_in_textframe(tf, mapping: list[tuple[str, str]]) -> None:
    for p in tf.paragraphs:
        full = p.text
        new = full
        for a, b in mapping:
            new = new.replace(a, b)
        if new != full:
            if p.runs:
                p.runs[0].text = new
                for r in p.runs[1:]:
                    r.text = ""
            else:
                p.text = new


def global_text_pass(slide, mapping: list[tuple[str, str]]) -> None:
    for shape in slide.shapes:
        if shape.has_text_frame:
            replace_in_textframe(shape.text_frame, mapping)


FINANCIAL_TITLE_MAP: list[tuple[str, str]] = [
    (
        "Financial model & projections",
        "Financial overview & optional forecasts",
    ),
    (
        "Source-backed pricing · Illustrative scenarios (finance to validate)\n"
        "List prices: product API (apps/api/src/routes/pricing.ts).",
        "Real subscription prices from the live product · Sample charts for discussion — "
        "finance replaces illustration data before external use.\n"
        "Prices match the product configuration (engineering path: pricing.ts).",
    ),
    ("Published SaaS list prices (monthly)", "Subscription list prices (monthly) — from the live product"),
    ("Published SaaS list prices (yearly contract)", "Subscription list prices (yearly) — from the live product"),
    ("Yearly list vs paying monthly × 12", "Savings: yearly plan vs twelve monthly payments (list price)"),
    (
        "Illustrative total scale ($M) — bear / base / bull",
        "Sample revenue paths ($M): cautious · middle · strong (illustration only)",
    ),
    (
        "Illustrative year-5 ARR by scenario ($M)",
        "Sample fifth-year recurring revenue by scenario ($M) — illustration",
    ),
    (
        "Illustrative revenue mix (year 5, base scenario)",
        "Sample revenue mix — fifth year, middle scenario (illustration)",
    ),
    (
        "Illustrative revenue build ($M ARR eq.) — stacked components",
        "Sample revenue build ($M) — stacked illustration (ARR-style view)",
    ),
    (
        "Illustrative EBITDA path ($M) — base scenario",
        "Sample operating result ($M, EBITDA-style) — middle scenario, illustration",
    ),
    (
        "Illustrative mature opex allocation (% of revenue)",
        "Sample operating cost mix (% of revenue) — illustration",
    ),
    (
        "Illustrative Y5 ARR sensitivity — blended ARPU shock",
        "Sample sensitivity — fifth-year revenue vs average revenue per account (illustration)",
    ),
    (
        "Illustrative MRR bridge ($k / mo) — template",
        "Sample monthly subscription revenue walkthrough — layout template (illustration)",
    ),
    (
        "Illustrative cumulative revenue index (base = 100 at Y1)",
        "Sample growth index (Year 1 = 100) — illustration, not actual performance",
    ),
    (
        "Financial charts — data sources & next steps",
        "Financial charts — what is real, what to replace",
    ),
]

GLOBAL_REPLACEMENTS: list[tuple[str, str]] = [
    ("VEX INVESTOR BRIEFING", "VEX — INVESTOR & PARTNER PACKAGE"),
    ("Institutional Investment Briefing", "Investment & partnership overview"),
    ("VEX Executive Investor Package", "VEX business & investor package"),
    ("Confidential", "Confidential — for qualified readers"),
]

FOOTER_REPLACEMENTS: list[tuple[str, str]] = [
    (
        "Illustrative projection — replace with finance / Stripe-backed figures before external use.",
        "Illustration only — replace with finance- and Stripe-backed figures before external use.",
    ),
    (
        "List prices: product API (apps/api/src/routes/pricing.ts).",
        "List prices match the live product catalog (see pricing.ts in source).",
    ),
    (
        "Shows list-price delta; actual invoice depends on Stripe Price IDs.",
        "Shows list-price comparison; actual charges follow your Stripe price IDs.",
    ),
]


def restyle_footer_textboxes(slide) -> None:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if not getattr(shape, "name", "").startswith("TextBox"):
            continue
        t = shape.text_frame.text.strip()
        if not t or len(t) > 220:
            continue
        if "Illustrative" in t or "List prices" in t or "Stripe" in t or "pricing.ts" in t:
            replace_in_textframe(shape.text_frame, FOOTER_REPLACEMENTS)
            style_textframe(shape.text_frame, footer=True)


def update_financial_bullet_slide(slide) -> None:
    """Financial appendix slide: friendlier bullets."""
    try:
        t = slide.shapes.title.text.strip()
    except ValueError:
        return
    if not t.startswith("Financial charts"):
        return
    for shape in slide.shapes:
        if not shape.has_text_frame or shape == slide.shapes.title:
            continue
        tf = shape.text_frame
        tf.clear()
        bullets = [
            "REAL (keep): Monthly $49 / $149 / $299 and yearly $470 / $1,430 / $2,870 — "
            "these match the default plan table in the live product (pricing API).",
            "ILLUSTRATION (refresh or delete): Curves, pies, stacked bars, EBITDA, sensitivity — "
            "drop in your finance model or remove before banks / press.",
            "LIVE DATA: Dealer metrics and pilot rollups come from the running platform "
            "(capital and pilot dashboards — ask your VEX contact).",
            "BEFORE YOU SEND: Finance signs off on every dollar; Stripe is truth if it differs from list.",
        ]
        tf.text = bullets[0]
        for line in bullets[1:]:
            p = tf.add_paragraph()
            p.text = line
            p.level = 0
        style_textframe(tf, title=False)


def rename_chart_series(chart) -> None:
    mapping = {"Bear": "Cautious", "Base": "Middle", "Bull": "Strong"}
    try:
        for plot in chart.plots:
            for series in plot.series:
                name = series.name
                if name in mapping:
                    series.name = mapping[name]
    except Exception:
        pass


def main() -> int:
    if not PPTX.is_file():
        print("Missing", PPTX, file=sys.stderr)
        return 1

    backup = Path(tempfile.gettempdir()) / "VEX-Pitch-Deck-theme-backup.pptx"
    shutil.copy2(PPTX, backup)

    prs = Presentation(PPTX)
    n = dedupe_financial_section(prs)
    if n:
        print("Removed duplicate slides:", n)

    # Slide master: dark base
    try:
        for master in prs.slide_masters:
            master.background.fill.solid()
            master.background.fill.fore_color.rgb = OBSIDIAN
    except Exception as e:
        print("Master bg note:", e)

    for slide in prs.slides:
        apply_slide_background(slide)
        recolor_auto_shapes(slide)
        global_text_pass(slide, GLOBAL_REPLACEMENTS)

        try:
            title_shape = slide.shapes.title
        except ValueError:
            title_shape = None

        for shape in slide.shapes:
            if shape.has_chart:
                style_chart(shape.chart)
                rename_chart_series(shape.chart)
            if not shape.has_text_frame:
                continue
            replace_in_textframe(shape.text_frame, FINANCIAL_TITLE_MAP)

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape.is_placeholder:
                ph = shape.placeholder_format
                if ph.type == PP_PLACEHOLDER.TITLE:
                    style_textframe(shape.text_frame, title=True)
                else:
                    style_textframe(shape.text_frame, title=False)
            else:
                if shape.top > 6_200_000:
                    style_textframe(shape.text_frame, footer=True)
                else:
                    style_textframe(shape.text_frame, title=False)

        restyle_footer_textboxes(slide)

    for slide in prs.slides:
        try:
            if slide.shapes.title and slide.shapes.title.text.strip().startswith("Financial charts"):
                update_financial_bullet_slide(slide)
                break
        except ValueError:
            continue

    prs.save(PPTX)
    print("Saved", PPTX, "slides:", len(prs.slides), "backup:", backup)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
